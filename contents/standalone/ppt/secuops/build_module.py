#!/usr/bin/env python3
"""
secuops module PPT generator — lecture md + lab yaml → pptx.

각 module 은 다음 구조:
1. Title slide (module 번호 + 제목 + 핵심 한 줄)
2. 학습 목표 (objectives)
3. 강의 시간 배분 표
4. 핵심 개념 (lecture md 의 ## 0~## N 헤더 별 1-2 slide)
5. 실습 step (lab.yaml 의 step 별 1 slide)
6. R/B/P 시나리오
7. 정리 + 다음 module 예고
"""

import sys
import re
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ─── 색상 팔레트 ──────────────────────────────────────────────
COLOR_TITLE_BG = RGBColor(0x1F, 0x3A, 0x68)  # 진한 남색
COLOR_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADING  = RGBColor(0x1F, 0x3A, 0x68)
COLOR_ACCENT   = RGBColor(0xE6, 0x4A, 0x19)  # 주황 (R/B/P 강조)
COLOR_CODE_BG  = RGBColor(0xF4, 0xF4, 0xF4)
COLOR_CODE_FG  = RGBColor(0x33, 0x33, 0x33)
COLOR_BODY     = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED    = RGBColor(0x66, 0x66, 0x66)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── slide helpers ─────────────────────────────────────────────
def add_title_slide(prs, module_no, title, subtitle):
    """Module 표지 slide."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 배경 사각형
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_TITLE_BG
    bg.line.fill.background()
    # module 번호
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.7))
    tf = tb.text_frame
    tf.text = f"Module {module_no:02d}"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xC1, 0x07)
    tf.paragraphs[0].font.bold = True
    # 제목
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_FG
    # 부제목
    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    # footer
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4))
    tf = tb.text_frame
    tf.text = "보안솔루션운영 (secops) · 6v6 학습 환경"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)


def add_section_header(prs, title, kicker=None):
    """섹션 구분 slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 좌측 색대
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.4), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_TITLE_BG
    bar.line.fill.background()
    if kicker:
        tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5))
        tf = tb.text_frame
        tf.text = kicker
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = COLOR_ACCENT
        tf.paragraphs[0].font.bold = True
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_HEADING


def add_content_slide(prs, title, bullets, footnote=None):
    """일반 본문 slide — 제목 + bullet list."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 제목
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING
    # 밑줄
    ln = s.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Emu(18000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = COLOR_HEADING
    ln.line.fill.background()
    # bullets
    body_top = Inches(1.5)
    body_h = Inches(5.5)
    tb = s.shapes.add_textbox(Inches(0.5), body_top, Inches(12.3), body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(b, tuple):
            level, text = b
        else:
            level, text = 0, b
        p.text = text
        p.level = level
        p.font.size = Pt(18 if level == 0 else 15)
        p.font.color.rgb = COLOR_BODY if level == 0 else COLOR_MUTED
        p.space_after = Pt(6)
    # footnote
    if footnote:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
        tf = tb.text_frame
        tf.text = footnote
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = COLOR_MUTED
        tf.paragraphs[0].font.italic = True


def add_table_slide(prs, title, headers, rows, footnote=None):
    """표 slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 제목
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING
    ln = s.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Emu(18000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = COLOR_HEADING
    ln.line.fill.background()
    # 표
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = s.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)).table
    # 헤더
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_TITLE_BG
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE_FG
    # 행
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_BODY
    if footnote:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
        tf = tb.text_frame
        tf.text = footnote
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = COLOR_MUTED


def add_code_slide(prs, title, code, caption=None):
    """code block slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING
    ln = s.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.3), Emu(18000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = COLOR_HEADING
    ln.line.fill.background()
    # code box
    box = s.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = 'Consolas'
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_CODE_FG
        p.space_after = Pt(2)
    if caption:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
        tf = tb.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLOR_MUTED
        tf.paragraphs[0].font.italic = True


def add_step_slide(prs, step_num, mission, commands, verify, why, target_vm):
    """lab step 1 slide — MISSION / 실행 명령 / 확인 / 왜."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 좌상단 step 번호
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(3.0), Inches(0.5))
    tf = tb.text_frame
    tf.text = f"STEP {step_num}"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    # 좌상단 target
    tb = s.shapes.add_textbox(Inches(9.0), Inches(0.25), Inches(3.8), Inches(0.5))
    tf = tb.text_frame
    tf.text = f"실행 대상: {target_vm}"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.RIGHT
    # MISSION
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = f"🎯 {mission}"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING
    # 실행 명령
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.4))
    tf = tb.text_frame
    tf.text = "💻 실행 명령"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_ACCENT
    box = s.shapes.add_shape(1, Inches(0.5), Inches(2.65), Inches(12.3), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.2)
    tf.text = commands
    for p in tf.paragraphs:
        p.font.name = 'Consolas'
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_CODE_FG
        p.space_after = Pt(1)
    # 확인 + 왜 (2 column)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(6.0), Inches(0.4))
    tf = tb.text_frame
    tf.text = "📋 확인 사항"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_ACCENT
    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(6.0), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = verify
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_BODY
    # 왜
    tb = s.shapes.add_textbox(Inches(7.0), Inches(5.3), Inches(5.8), Inches(0.4))
    tf = tb.text_frame
    tf.text = "💡 왜 이걸 하는가?"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_ACCENT
    tb = s.shapes.add_textbox(Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = why
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_BODY


def add_thanks_slide(prs, next_module=None):
    """마지막 정리 slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_TITLE_BG
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.text = "본 module 종료"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_TITLE_FG
    if next_module:
        tb = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = f"다음 module → {next_module}"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xC1, 0x07)


# ─── md / yaml 파서 ────────────────────────────────────────────
def parse_lecture_md(md_path):
    """lecture md → {title, objectives, chapters: [{title, bullets}]}."""
    s = Path(md_path).read_text()
    title_m = re.match(r"# Week \d+ — (.+)", s.split("\n")[0])
    title = title_m.group(1) if title_m else Path(md_path).stem
    out = {"title": title, "objectives": [], "chapters": []}
    # objectives
    obj_match = re.search(r"## 학습 목표\n(.*?)\n---", s, re.DOTALL)
    if obj_match:
        for line in obj_match.group(1).strip().split("\n"):
            ln = line.strip()
            if ln and (ln[0].isdigit() or ln.startswith("-")):
                ln = re.sub(r"^\d+\.\s*", "", ln)
                ln = re.sub(r"^-\s*", "", ln)
                out["objectives"].append(ln.replace("**", ""))
    # chapters (## N. xxx ~ next ##)
    chap_re = re.findall(r"\n## (\d+\..+?)\n(.*?)(?=\n## \d+\.|\n## \S|\Z)", s, re.DOTALL)
    for title, body in chap_re:
        bullets = []
        for line in body.split("\n"):
            ln = line.rstrip()
            if not ln:
                continue
            if ln.startswith("###"):
                bullets.append((0, ln.lstrip("# ").strip().replace("**", "")))
            elif ln.startswith("- ") or ln.startswith("* "):
                txt = ln[2:].strip().replace("**", "")
                bullets.append((1, txt[:140]))
            elif ln.startswith("  - "):
                bullets.append((2, ln[4:].strip().replace("**", "")[:140]))
        if bullets:
            out["chapters"].append({"title": title.strip(), "bullets": bullets[:14]})
    return out


def parse_lab_yaml(yaml_path):
    """lab.yaml → [{order, mission, commands, verify, why, target}]."""
    d = yaml.safe_load(Path(yaml_path).read_text())
    steps = []
    for s in d.get("steps", []):
        instr = s.get("instruction", "") or ""
        mission = ""
        commands = ""
        verify_txt = ""
        why = ""
        if "🎯 MISSION" in instr:
            m = re.search(r"🎯 MISSION\s*\n+(.*?)(?=\n##|\Z)", instr, re.DOTALL)
            if m:
                mission = m.group(1).strip().split("\n\n")[0][:200]
        if "💻 실행 명령" in instr:
            m = re.search(r"💻 실행 명령\s*\n+(?:```bash\s*\n)?(.*?)(?=\n##|```)", instr, re.DOTALL)
            if m:
                commands = m.group(1).strip()[:1200]
        if "📋 확인 사항" in instr:
            m = re.search(r"📋 확인 사항\s*\n+(.*?)(?=\n##|\Z)", instr, re.DOTALL)
            if m:
                verify_txt = m.group(1).strip()[:400]
        if "💡 왜" in instr:
            m = re.search(r"💡 왜.*?\n+(.*?)(?=\n##|\Z)", instr, re.DOTALL)
            if m:
                why = m.group(1).strip()[:300]
        if not mission:
            mission = instr.split("\n")[0][:200] if instr else f"Step {s.get('order')}"
        steps.append({
            "order": s.get("order"),
            "mission": mission,
            "commands": commands or "(상세 명령은 lab YAML 참조)",
            "verify": verify_txt or "(verify 절차는 lab YAML 참조)",
            "why": why or "(학습 의의는 lab YAML 참조)",
            "target": s.get("target_vm", "?"),
        })
    return {"title": d.get("title", Path(yaml_path).stem), "steps": steps}


# ─── module 생성 main ─────────────────────────────────────────
def build_module(week_no, out_path):
    md = f"contents/standalone/lecture/secuops/week{week_no:02d}.md"
    yml = f"contents/standalone/lab/secuops/week{week_no:02d}.yaml"
    lec = parse_lecture_md(md)
    lab = parse_lab_yaml(yml)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    add_title_slide(prs, week_no, lec["title"], lab["title"])

    # 2. 학습 목표
    if lec["objectives"]:
        add_content_slide(prs, "학습 목표", [(0, o) for o in lec["objectives"]],
                          footnote=f"보안솔루션운영 (secops) · Module {week_no:02d}")

    # 3. 강의 + 실습 개요 (section)
    add_section_header(prs, lec["title"], kicker=f"Module {week_no:02d} — 이론 + 실습")

    # 4. 핵심 개념 (chapters)
    for ch in lec["chapters"][:8]:
        bullets = ch["bullets"][:10]
        if not bullets:
            continue
        add_content_slide(prs, ch["title"], bullets)

    # 5. 실습 section
    add_section_header(prs, "실습 (Hands-on Lab)", kicker=f"{len(lab['steps'])} STEP")

    # 6. lab steps
    for st in lab["steps"]:
        add_step_slide(prs, st["order"], st["mission"], st["commands"],
                       st["verify"], st["why"], st["target"])

    # 7. 정리
    next_no = week_no + 1
    next_label = f"Module {next_no:02d}" if next_no <= 15 else "과정 종료"
    add_thanks_slide(prs, next_module=next_label)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"✅ {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        for arg in sys.argv[1:]:
            no = int(arg)
            build_module(no, f"contents/standalone/ppt/secuops/module{no:02d}.pptx")
    else:
        for no in range(1, 16):
            build_module(no, f"contents/standalone/ppt/secuops/module{no:02d}.pptx")
